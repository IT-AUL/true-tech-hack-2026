import re

with open('/Users/dremotha/Projects/Temp/true-tech-hack-2026/backend/open_webui/tools/builtin.py') as f:
    text = f.read()


def replacer(match):
    return """async def generate_presentation(
    topic: str,
    num_slides: int = 5,
    __request__: Request = None,
    __user__: dict = None,
    __event_emitter__: callable = None,
    __chat_id__: str = None,
    __message_id__: str = None,
    __model__: dict = None,
) -> str:
    \"\"\"
    Generate a PPTX presentation about a specific topic. Use this tool autonomously when the user asks for a presentation.

    :param topic: Detailed topic or context for the presentation (can include web research results)
    :param num_slides: Desired number of slides (default: 5)
    :return: JSON result with file linking data
    \"\"\"
    if __request__ is None:
        return json.dumps({'error': 'Request context not available'})
    if not __user__:
        return json.dumps({'error': 'User context not available'})

    try:
        import asyncio
        import io
        import uuid
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.util import Pt
        from open_webui.models.files import FileForm, Files
        from open_webui.storage.provider import Storage
        from open_webui.models.users import UserModel
        from open_webui.utils.chat import generate_chat_completion

        orchestrator_model = __model__["id"] if isinstance(__model__, dict) and "id" in __model__ else "default"

        if __event_emitter__:
            await __event_emitter__({
                "type": "status",
                "data": {"action": "presentation", "description": "Создание структуры презентации (Планировщик)...", "done": False}
            })

        async def _safe_generate(payload, retries=3):
            user_model = UserModel(**__user__)
            for attempt in range(retries):
                try:
                    return await generate_chat_completion(__request__, form_data=payload, user=user_model)
                except Exception as e:
                    if attempt == retries - 1: raise e
                    await asyncio.sleep(2)
            return None

        # --- MAP REDUCE LOGIC ---
        # 1. PLANNER
        generation_prompt = f\"\"\"Ты Планировщик - эксперт по структуре презентаций.
Пользователю нужна презентация на тему: "{topic}" (Желаемое количество слайдов: {num_slides})
Твоя задача написать ТОЛЬКО базовую структуру (список слайдов и короткую суть каждого слайда).

Ответь ИСКЛЮЧИТЕЛЬНО валидным JSON объектом следующей структуры:
{{
  "title": "Название презентации",
  "theme": "light",
  "slides": [
    {{
      "title": "Заголовок слайда",
      "description": "О чем должен быть этот слайд (1-2 предложения)"
    }}
  ]
}}
\"\"\"
        gen_payload = {
            "model": orchestrator_model,
            "messages": [{"role": "user", "content": generation_prompt}],
            "stream": False,
        }

        gen_resp = await _safe_generate(gen_payload)
        gen_dict = json.loads(gen_resp.body.decode('utf-8')) if hasattr(gen_resp, 'body') else gen_resp
        content = gen_dict.get("choices", [{}])[0].get("message", {}).get("content", "")
        content = content[content.find('{') : content.rfind('}') + 1]
        parsed_pres = json.loads(content)

        outline_slides = parsed_pres.get("slides", [])
        pres_title = parsed_pres.get("title", "Презентация")
        theme = parsed_pres.get("theme", "light")
        slides_json_data = []

        # 2. WORKER
        for i, slide_plan in enumerate(outline_slides[:num_slides]):
            if __event_emitter__:
                await __event_emitter__({
                    "type": "status",
                    "data": {"action": "presentation", "description": f"Генерация контента слайда {i+1} из {len(outline_slides[:num_slides])}...", "done": False}
                })
            slide_prompt = f\"\"\"Ты Писатель. Создай контент для ОДНОГО слайда в рамках презентации "{pres_title}".
Текущий слайд: "{slide_plan.get('title')}"
Задача: {slide_plan.get('description')}

Твоя задача — написать буллеты (тезисы) и текст выступления (заметки для спикера). Пиши на русском языке.
Ответь ИСКЛЮЧИТЕЛЬНО валидным JSON объектом:
{{
  "title": "{slide_plan.get('title')}",
  "bullets": ["Подробный тезис 1", "Подробный тезис 2", "Подробный тезис 3"],
  "notes": "Детальный текст для спикера к этому слайду."
}}
\"\"\"
            slide_payload = {
                "model": orchestrator_model,
                "messages": [{"role": "user", "content": slide_prompt}],
                "stream": False,
            }
            slide_resp = await _safe_generate(slide_payload)
            slide_dict = json.loads(slide_resp.body.decode('utf-8')) if hasattr(slide_resp, 'body') else slide_resp
            slide_content = slide_dict.get("choices", [{}])[0].get("message", {}).get("content", "")
            slide_content = slide_content[slide_content.find('{') : slide_content.rfind('}') + 1]
            try:
                slides_json_data.append(json.loads(slide_content))
            except:
                slides_json_data.append({"title": slide_plan.get("title", "Слайд"), "bullets": ["Ошибка"], "notes": ""})

        # --- PPTX GENERATION ---
        if __event_emitter__:
            await __event_emitter__({
                "type": "status",
                "data": {"action": "presentation", "description": "Сборка файла PPTX...", "done": False}
            })

        prs = Presentation()

        def apply_theme_text(paragraph, selected_theme: str) -> None:
            for run in paragraph.runs:
                if selected_theme == 'dark':
                    run.font.color.rgb = RGBColor(255, 255, 255)
                elif selected_theme == 'corporate':
                    run.font.color.rgb = RGBColor(25, 55, 109)
                else:
                    run.font.color.rgb = RGBColor(0, 0, 0)

        title_layout = prs.slide_layouts[0]
        first_slide = prs.slides.add_slide(title_layout)
        first_slide.shapes.title.text = pres_title
        subtitle = first_slide.placeholders[1]
        subtitle.text = 'Generated by Agentic Pipeline'
        apply_theme_text(first_slide.shapes.title.text_frame.paragraphs[0], theme)
        apply_theme_text(subtitle.text_frame.paragraphs[0], theme)

        content_layout = prs.slide_layouts[1]
        for slide_data in slides_json_data:
            slide_title = str(slide_data.get('title') or 'Slide')
            bullets = slide_data.get('bullets') or []
            notes = str(slide_data.get('notes') or '')

            slide = prs.slides.add_slide(content_layout)
            slide.shapes.title.text = slide_title[:150]
            apply_theme_text(slide.shapes.title.text_frame.paragraphs[0], theme)

            text_frame = slide.placeholders[1].text_frame
            text_frame.clear()

            if isinstance(bullets, list) and bullets:
                for idx, bullet in enumerate(bullets[:12]):
                    paragraph = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
                    paragraph.text = str(bullet)[:300]
                    paragraph.level = 0
                    paragraph.font.size = Pt(22 if idx == 0 else 18)
                    apply_theme_text(paragraph, theme)
            else:
                paragraph = text_frame.paragraphs[0]
                paragraph.text = 'No bullet points provided.'
                paragraph.font.size = Pt(18)
                apply_theme_text(paragraph, theme)

            if notes:
                try:
                    slide.notes_slide.notes_text_frame.text = notes[:2000]
                except Exception:
                    pass

        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        file_bytes = output.getvalue()

        file_id = str(uuid.uuid4())
        cleaned = ''.join(c if c.isalnum() or c in ('-', '_', ' ') else '_' for c in pres_title).strip()
        cleaned = cleaned.replace(' ', '_') or 'presentation'
        filename = f'{cleaned}.pptx'
        storage_filename = f'{file_id}_{filename}'

        user = UserModel(**__user__)
        _, file_path = Storage.upload_file(
            io.BytesIO(file_bytes),
            storage_filename,
            {
                'OpenWebUI-User-Email': user.email,
                'OpenWebUI-User-Id': user.id,
                'OpenWebUI-User-Name': user.name,
                'OpenWebUI-File-Id': file_id,
            },
        )

        file_item = Files.insert_new_file(
            user.id,
            FileForm(
                id=file_id,
                filename=filename,
                path=file_path,
                data={'content': ''},
                meta={
                    'name': filename,
                    'content_type': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
                    'size': len(file_bytes),
                },
            ),
        )

        file_payload = [
            {
                'type': 'file',
                'id': file_item.id,
                'url': file_item.id,
                'name': file_item.filename,
                'content_type': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
                'size': len(file_bytes),
                'status': 'uploaded',
            }
        ]

        if __chat_id__ and __message_id__:
            db_files = Chats.add_message_files_by_id_and_message_id(__chat_id__, __message_id__, file_payload)
            if db_files is not None:
                file_payload = db_files

        if __event_emitter__:
            await __event_emitter__(
                {
                    'type': 'chat:message:files',
                    'data': {
                        'files': file_payload,
                    },
                }
            )

        if __event_emitter__:
            await __event_emitter__({
                "type": "status",
                "data": {"action": "presentation", "description": "Презентация готова!", "done": True}
            })

        return f"Презентация была успешно сгенерирована. Файл называется {filename}. Уведоми пользователя и скажи ему скачать файл."
    except Exception as e:
        log.exception(f'generate_presentation error: {e}')
        return f"Произошла ошибка при генерации презентации: {e}"
"""


pattern = re.compile(r'async def generate_presentation\(.*?return json\.dumps\(\{\'error\': str\(e\)\}\)', re.DOTALL)
new_text = pattern.sub(replacer, text)

with open('/Users/dremotha/Projects/Temp/true-tech-hack-2026/backend/open_webui/tools/builtin.py', 'w') as f:
    f.write(new_text)
