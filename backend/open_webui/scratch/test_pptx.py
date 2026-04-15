import io

from pptx import Presentation


def test_pptx():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = 'Test Title'

    output = io.BytesIO()
    prs.save(output)
    size = len(output.getvalue())
    print(f'PPTX Size: {size} bytes')

    # Check if it's a valid ZIP
    import zipfile

    output.seek(0)
    try:
        with zipfile.ZipFile(output) as zf:
            print(f'ZIP valid. Contains {len(zf.namelist())} files.')
    except Exception as e:
        print(f'ZIP INVALID: {e}')


if __name__ == '__main__':
    test_pptx()
