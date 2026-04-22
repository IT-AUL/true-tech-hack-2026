"""
Presentation pipeline — agentic backend orchestrator that turns a user prompt
into a downloadable .pptx file.

Stages (all async, can emit progress via an event_emitter):
  1. build_outline()     — single LLM call returns the full slide plan as JSON
  2. detail_slides()     — (optional) per-slide refinement in parallel
  3. generate_slide_images() — optional image_gen per slide, embedded into pptx
  4. assemble_pptx()     — python-pptx → bytes

The top-level entry point is run_presentation_pipeline(), invoked from
routers/openai.py::_auto_routed_presentation_attempt when auto-routing
resolves category == 'presentation'.
"""

from __future__ import annotations

import asyncio
import io
import json
import logging
import os
import re
from dataclasses import asdict, dataclass, field
from typing import Any

log = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# Data classes
# ---------------------------------------------------------------------------


@dataclass(slots=True)
class SlideSpec:
    title: str
    bullets: list[str] = field(default_factory=list)
    speaker_notes: str = ''
    # LLM sets this when an illustration would genuinely help the slide
    image_prompt: str | None = None


@dataclass(slots=True)
class PresentationSpec:
    title: str
    subtitle: str | None
    slides: list[SlideSpec]
    language: str = 'ru'
    theme: str = 'default'


# ---------------------------------------------------------------------------
# LLM wrapper — reuses the active OpenAI-compatible provider of the chosen writer model
# ---------------------------------------------------------------------------


def _resolve_writer_endpoint(request, writer_model_id: str) -> tuple[str, str]:
    """Return (base_url, api_key) for the writer model based on its urlIdx in OPENAI_MODELS."""
    models = getattr(getattr(request, 'app', None), 'state', None)
    models = getattr(models, 'OPENAI_MODELS', None) or {}
    meta = models.get(writer_model_id, {}) if isinstance(models, dict) else {}
    idx = meta.get('urlIdx', 0) if isinstance(meta, dict) else 0
    try:
        base_url = request.app.state.config.OPENAI_API_BASE_URLS[idx]
        api_key = request.app.state.config.OPENAI_API_KEYS[idx]
        return base_url, api_key
    except (AttributeError, IndexError, KeyError):
        base_url = os.environ.get('OPENAI_API_BASE_URL', 'https://api.openai.com/v1')
        api_key = os.environ.get('OPENAI_API_KEY', '')
        return base_url, api_key


async def _llm_json(
    request,
    writer_model_id: str,
    messages: list[dict[str, Any]],
    *,
    max_tokens: int = 4096,
    timeout: float = 120.0,
    temperature: float = 0.4,
) -> dict[str, Any]:
    """
    Call the writer model expecting a JSON object in the response. Falls back
    to best-effort JSON extraction (regex) if the provider does not honour
    response_format.
    """
    import httpx  # lazy import so assemble_pptx can be used without httpx installed

    base_url, api_key = _resolve_writer_endpoint(request, writer_model_id)

    body: dict[str, Any] = {
        'model': writer_model_id,
        'messages': messages,
        'max_tokens': max_tokens,
        'temperature': temperature,
        'stream': False,
        'response_format': {'type': 'json_object'},
    }
    headers = {
        'Authorization': f'Bearer {api_key}',
        'Content-Type': 'application/json',
    }

    async with httpx.AsyncClient(timeout=timeout) as client:
        try:
            resp = await client.post(f'{base_url}/chat/completions', headers=headers, json=body)
            resp.raise_for_status()
        except httpx.HTTPStatusError as exc:
            # Some providers reject response_format — retry without it.
            if exc.response is not None and exc.response.status_code in (400, 422):
                body.pop('response_format', None)
                resp = await client.post(f'{base_url}/chat/completions', headers=headers, json=body)
                resp.raise_for_status()
            else:
                raise

    data = resp.json()
    raw = (data.get('choices', [{}])[0].get('message', {}).get('content', '') or '').strip()

    # Strip possible ```json fences
    if raw.startswith('```'):
        raw = re.sub(r'^```[a-zA-Z]*\n?', '', raw)
        raw = re.sub(r'\n?```$', '', raw)

    try:
        return json.loads(raw)
    except json.JSONDecodeError:
        match = re.search(r'\{.*\}', raw, re.DOTALL)
        if match:
            try:
                return json.loads(match.group(0))
            except json.JSONDecodeError:
                pass
        log.warning('presentation pipeline: failed to parse JSON from writer (first 400 chars): %s', raw[:400])
        raise ValueError('Writer model did not return valid JSON')


# ---------------------------------------------------------------------------
# Stage 1: outline
# ---------------------------------------------------------------------------


_OUTLINE_SYSTEM_PROMPT = """\
Ты — профессиональный автор презентаций. На вход ты получаешь запрос пользователя и возвращаешь полностью готовый план презентации в виде строгого JSON.

Правила:
- Минимум 5 и максимум {max_slides} слайдов (включая титульный).
- Первый слайд всегда title-слайд с subtitle (краткий подзаголовок).
- Далее — content-слайды: короткий title, 3–5 bullets по 5–12 слов каждый.
- speaker_notes — 2–4 предложения, что сказать вслух на этом слайде.
- Если иллюстрация реально помогает (концепт, метафора, схема, продукт) — добавь image_prompt (ёмкое описание для image-gen, 1–2 предложения, англ. или рус.). Для большинства слайдов image_prompt = null.
- Пиши на языке запроса пользователя.
- Контент должен быть конкретным, без воды и без повторов соседних слайдов.

Формат ответа — ТОЛЬКО JSON объект:
{{
  "title": "Заголовок презентации",
  "subtitle": "Подзаголовок или null",
  "language": "ru" | "en",
  "slides": [
    {{"title": "...", "bullets": ["...", "..."], "speaker_notes": "...", "image_prompt": null}},
    ...
  ]
}}
"""


async def build_outline(
    request,
    writer_model_id: str,
    user_prompt: str,
    *,
    max_slides: int = 12,
    extra_context: str | None = None,
) -> PresentationSpec:
    system_prompt = _OUTLINE_SYSTEM_PROMPT.format(max_slides=max_slides)
    user_content = user_prompt.strip()
    if extra_context:
        user_content = f'{user_content}\n\nДополнительный контекст:\n{extra_context.strip()}'

    data = await _llm_json(
        request,
        writer_model_id,
        [
            {'role': 'system', 'content': system_prompt},
            {'role': 'user', 'content': user_content},
        ],
        max_tokens=4096,
    )
    return _parse_outline(data, max_slides=max_slides)


def _parse_outline(data: dict[str, Any], *, max_slides: int) -> PresentationSpec:
    if not isinstance(data, dict):
        raise ValueError('Outline JSON is not an object')

    title = str(data.get('title') or 'Презентация').strip() or 'Презентация'
    subtitle_raw = data.get('subtitle')
    subtitle = str(subtitle_raw).strip() if subtitle_raw else None
    language = (str(data.get('language') or 'ru') or 'ru').strip().lower()
    if language not in ('ru', 'en'):
        language = 'ru'

    raw_slides = data.get('slides') or []
    if not isinstance(raw_slides, list) or not raw_slides:
        raise ValueError('Outline JSON has no slides')

    slides: list[SlideSpec] = []
    for raw in raw_slides[:max_slides]:
        if not isinstance(raw, dict):
            continue
        s_title = str(raw.get('title') or '').strip()
        if not s_title:
            continue
        bullets = raw.get('bullets') or []
        if not isinstance(bullets, list):
            bullets = []
        clean_bullets = [str(b).strip() for b in bullets if str(b).strip()][:6]
        notes = str(raw.get('speaker_notes') or '').strip()
        img_prompt_raw = raw.get('image_prompt')
        img_prompt = str(img_prompt_raw).strip() if img_prompt_raw else None
        if img_prompt and img_prompt.lower() in ('null', 'none', ''):
            img_prompt = None
        slides.append(SlideSpec(title=s_title, bullets=clean_bullets, speaker_notes=notes, image_prompt=img_prompt))

    if not slides:
        raise ValueError('Outline JSON produced zero valid slides')

    return PresentationSpec(title=title, subtitle=subtitle, slides=slides, language=language)


# ---------------------------------------------------------------------------
# Stage 2: per-slide detailing (parallel, optional)
# ---------------------------------------------------------------------------


_DETAIL_SYSTEM_PROMPT = """\
Ты — редактор слайдов. Твоя задача — улучшить ОДИН слайд презентации: сделать bullets более конкретными и информативными, переписать speaker_notes, сохранить тему и тон.

Правила:
- Сохрани title (можно лишь слегка уточнить).
- 3–5 bullets, каждое 5–14 слов, конкретика, без пустых общих фраз.
- speaker_notes: 2–4 предложения, то, что ведущий скажет вслух.
- Не выдумывай цифры и факты, если их нет в контексте.
- Сохрани image_prompt, если он был задан в оригинальном слайде.

Ответ — ТОЛЬКО JSON объект слайда:
{"title": "...", "bullets": ["..."], "speaker_notes": "...", "image_prompt": null | "..."}\
"""


async def _detail_one_slide(
    request,
    writer_model_id: str,
    spec: PresentationSpec,
    idx: int,
) -> SlideSpec:
    slide = spec.slides[idx]
    neighbors = []
    if idx > 0:
        neighbors.append(f'Предыдущий слайд: {spec.slides[idx - 1].title}')
    if idx < len(spec.slides) - 1:
        neighbors.append(f'Следующий слайд: {spec.slides[idx + 1].title}')

    context = (
        f'Тема презентации: {spec.title}\n'
        + (f'Подзаголовок: {spec.subtitle}\n' if spec.subtitle else '')
        + f'Язык: {spec.language}\n'
        + (('\n'.join(neighbors) + '\n') if neighbors else '')
        + f'\nТекущий слайд ({idx + 1}/{len(spec.slides)}):\n{json.dumps(asdict(slide), ensure_ascii=False, indent=2)}'
    )

    try:
        data = await _llm_json(
            request,
            writer_model_id,
            [
                {'role': 'system', 'content': _DETAIL_SYSTEM_PROMPT},
                {'role': 'user', 'content': context},
            ],
            max_tokens=1024,
            temperature=0.5,
        )
    except Exception as exc:
        log.warning('detail_slide %d failed, keeping outline version: %s', idx, exc)
        return slide

    title = str(data.get('title') or slide.title).strip() or slide.title
    bullets_raw = data.get('bullets') or slide.bullets
    bullets = (
        [str(b).strip() for b in bullets_raw if str(b).strip()][:6] if isinstance(bullets_raw, list) else slide.bullets
    )
    notes = str(data.get('speaker_notes') or slide.speaker_notes).strip()
    img_prompt_raw = data.get('image_prompt')
    img_prompt = str(img_prompt_raw).strip() if img_prompt_raw else slide.image_prompt
    if img_prompt and img_prompt.lower() in ('null', 'none', ''):
        img_prompt = None

    return SlideSpec(title=title, bullets=bullets, speaker_notes=notes, image_prompt=img_prompt)


async def detail_slides(
    request,
    writer_model_id: str,
    spec: PresentationSpec,
    *,
    skip_title_slide: bool = True,
    on_progress=None,
) -> PresentationSpec:
    """
    Refine every (non-title) slide in parallel. Failures fall back to the outline version.
    on_progress(done, total) is called after each slide completes.
    """
    total = len(spec.slides)
    done_counter = {'n': 0}

    async def runner(idx: int) -> SlideSpec:
        if skip_title_slide and idx == 0:
            result = spec.slides[idx]
        else:
            result = await _detail_one_slide(request, writer_model_id, spec, idx)
        done_counter['n'] += 1
        if on_progress is not None:
            try:
                await on_progress(done_counter['n'], total)
            except Exception as exc:
                log.debug('detail_slides on_progress callback failed: %s', exc)
        return result

    results = await asyncio.gather(*[runner(i) for i in range(total)], return_exceptions=False)
    spec.slides = list(results)
    return spec


# ---------------------------------------------------------------------------
# Stage 3: per-slide image generation (parallel, bounded, optional)
# ---------------------------------------------------------------------------


_FILE_URL_ID_RE = re.compile(r'/files/([0-9a-f\-]+)/content', re.IGNORECASE)


async def _generate_one_slide_image(
    request,
    user,
    image_model_id: str | None,
    prompt: str,
    sem: asyncio.Semaphore,
    *,
    timeout: float = 120.0,
) -> bytes | None:
    """Call image_generations without chat attachment; return raw image bytes or None."""
    from open_webui.models.files import Files
    from open_webui.routers.images import CreateImageForm, image_generations
    from open_webui.storage.provider import Storage

    async with sem:
        try:
            images = await asyncio.wait_for(
                image_generations(
                    request=request,
                    form_data=CreateImageForm(model=image_model_id or '', prompt=prompt),
                    metadata={},  # no chat_id/message_id → do NOT attach to chat
                    user=user,
                ),
                timeout=timeout,
            )
        except Exception as exc:
            log.warning('presentation image_gen failed for prompt %r: %s', prompt[:80], exc)
            return None

    if not images:
        return None

    url = images[0].get('url') if isinstance(images[0], dict) else None
    if not url:
        return None

    match = _FILE_URL_ID_RE.search(url)
    if not match:
        return None
    file_id = match.group(1)

    try:
        file_item = Files.get_file_by_id(file_id)
        if not file_item or not file_item.path:
            return None
        local_path = Storage.get_file(file_item.path)
        with open(local_path, 'rb') as fh:
            data = fh.read()
        # Best-effort cleanup; the image lives in storage otherwise forever.
        try:
            Files.delete_file_by_id(file_id)
        except Exception:
            pass
        try:
            Storage.delete_file(file_item.path)
        except Exception:
            pass
        return data
    except Exception as exc:
        log.warning('presentation failed to read image bytes for file_id=%s: %s', file_id, exc)
        return None


def _pick_image_model(request) -> str | None:
    """Return the first image-gen-capable model id from the catalog, or None."""
    models = getattr(getattr(request, 'app', None), 'state', None)
    models = getattr(models, 'OPENAI_MODELS', None) or {}
    if not isinstance(models, dict):
        return None

    # Prefer flagship endpoints.
    preferred_order = ('gpt-image', 'dall-e-3', 'imagen', 'flux-pro', 'flux.1', 'flux', 'qwen-image', 'sdxl', 'sd3')
    entries = [(mid, m) for mid, m in models.items() if isinstance(m, dict) and mid != 'auto']

    def _search_text(m: dict) -> str:
        parts = [m.get('id', ''), m.get('name', ''), m.get('owned_by', '')]
        return ' '.join(str(p) for p in parts if p).lower()

    for pref in preferred_order:
        for mid, m in entries:
            if pref in _search_text(m):
                return mid
    return None


async def generate_slide_images(
    request,
    user,
    spec: PresentationSpec,
    *,
    concurrency: int = 3,
    on_progress=None,
) -> dict[int, bytes]:
    """Generate images for slides where image_prompt is set. Returns {slide_idx: bytes}."""
    targets = [(i, s.image_prompt) for i, s in enumerate(spec.slides) if s.image_prompt]
    if not targets:
        return {}

    image_model = _pick_image_model(request)
    if image_model is None:
        log.info('presentation: no image-gen model available, skipping illustrations')
        return {}

    sem = asyncio.Semaphore(max(1, concurrency))
    total = len(targets)
    done = {'n': 0}

    async def run_one(idx: int, prompt: str) -> tuple[int, bytes | None]:
        data = await _generate_one_slide_image(request, user, image_model, prompt, sem)
        done['n'] += 1
        if on_progress is not None:
            try:
                await on_progress(done['n'], total)
            except Exception as exc:
                log.debug('image on_progress callback failed: %s', exc)
        return idx, data

    results = await asyncio.gather(*[run_one(i, p) for i, p in targets])
    return {idx: data for idx, data in results if data}


# ---------------------------------------------------------------------------
# Stage 4: assemble .pptx
# ---------------------------------------------------------------------------


def _set_ph_geometry(ph, *, left_in: float, top_in: float, width_in: float, height_in: float) -> None:
    """Force explicit geometry on a placeholder.

    python-pptx lets placeholders inherit position from the slide layout and then
    omits <a:xfrm> in the per-slide XML. That breaks lightweight downstream
    renderers (our frontend pptxToImages skips any shape without xfrm), so we
    always write an explicit absolute transform on every placeholder we touch.
    """
    from pptx.util import Inches

    ph.left = Inches(left_in)
    ph.top = Inches(top_in)
    ph.width = Inches(width_in)
    ph.height = Inches(height_in)


def assemble_pptx(spec: PresentationSpec, images: dict[int, bytes] | None = None) -> bytes:
    """Build a minimal, clean PPTX in memory. Returns raw bytes."""
    # Local import so module import cost stays low when python-pptx is unused.
    from pptx import Presentation
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    images = images or {}
    prs = Presentation()
    # 16:9 canvas (default python-pptx is 4:3 10x7.5)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_w_in = 13.333

    # --- Title slide ---
    title_layout = prs.slide_layouts[0]  # Title Slide
    title_slide = prs.slides.add_slide(title_layout)

    if title_slide.shapes.title is not None:
        title_shape = title_slide.shapes.title
        title_shape.text = spec.title
        _set_ph_geometry(
            title_shape,
            left_in=0.8,
            top_in=2.4,
            width_in=slide_w_in - 1.6,
            height_in=1.8,
        )
        for p in title_shape.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.size = Pt(40)
                run.font.bold = True

    if spec.subtitle:
        for ph in title_slide.placeholders:
            if ph.placeholder_format.idx == 1:
                ph.text = spec.subtitle
                _set_ph_geometry(
                    ph,
                    left_in=0.8,
                    top_in=4.4,
                    width_in=slide_w_in - 1.6,
                    height_in=1.0,
                )
                for p in ph.text_frame.paragraphs:
                    p.alignment = PP_ALIGN.CENTER
                    for run in p.runs:
                        run.font.size = Pt(22)
                break

    # Notes on title slide (notes_slide is lazily created on access).
    try:
        title_slide.notes_slide.notes_text_frame.text = spec.slides[0].speaker_notes or ''
    except Exception as exc:
        log.debug('Failed to set title-slide notes: %s', exc)

    # --- Content slides ---
    content_layout = prs.slide_layouts[1]  # Title and Content

    for i, slide in enumerate(spec.slides):
        if i == 0:
            # The title slide content has already been placed via the first slide.
            continue
        slide_obj = prs.slides.add_slide(content_layout)

        has_image = bool(images.get(i))

        if slide_obj.shapes.title is not None:
            title_shape = slide_obj.shapes.title
            title_shape.text = slide.title
            _set_ph_geometry(
                title_shape,
                left_in=0.5,
                top_in=0.35,
                width_in=slide_w_in - 1.0,
                height_in=1.0,
            )
            for p in title_shape.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(28)
                    run.font.bold = True

        # Find the content placeholder (idx 1)
        body_ph = None
        for ph in slide_obj.placeholders:
            if ph.placeholder_format.idx == 1:
                body_ph = ph
                break

        if body_ph is not None:
            if has_image:
                _set_ph_geometry(body_ph, left_in=0.5, top_in=1.55, width_in=6.3, height_in=5.5)
            else:
                _set_ph_geometry(
                    body_ph,
                    left_in=0.5,
                    top_in=1.55,
                    width_in=slide_w_in - 1.0,
                    height_in=5.5,
                )

            if slide.bullets:
                tf = body_ph.text_frame
                tf.word_wrap = True
                tf.clear()
                for j, bullet in enumerate(slide.bullets):
                    p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                    p.text = bullet
                    p.level = 0
                    for run in p.runs:
                        run.font.size = Pt(20)

        if has_image:
            try:
                slide_obj.shapes.add_picture(
                    io.BytesIO(images[i]),
                    left=Inches(7.1),
                    top=Inches(1.55),
                    width=Inches(5.7),
                    height=Inches(5.5),
                )
            except Exception as exc:
                log.warning('Failed to embed slide image on slide %d: %s', i, exc)

        # Speaker notes — notes_slide is lazily created on access.
        if slide.speaker_notes:
            try:
                slide_obj.notes_slide.notes_text_frame.text = slide.speaker_notes
            except Exception as exc:
                log.debug('Failed to set slide %d notes: %s', i, exc)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# Top-level orchestrator
# ---------------------------------------------------------------------------


async def run_presentation_pipeline(
    request,
    user,
    *,
    user_prompt: str,
    writer_model_id: str,
    max_slides: int = 12,
    enable_images: bool = True,
    enable_detail_step: bool = True,
    image_concurrency: int = 3,
    extra_context: str | None = None,
    event_emitter=None,
) -> tuple[PresentationSpec, bytes]:
    """
    Run the full pipeline and return (spec, pptx_bytes).

    event_emitter is an async callable (same shape as sockets' event_emitter)
    used to stream progress status updates to the UI.
    """

    async def status(description: str, *, done: bool = False, extra: dict | None = None):
        if event_emitter is None:
            return
        data = {
            'action': 'presentation_progress',
            'description': description,
            'done': done,
        }
        if extra:
            data.update(extra)
        try:
            await event_emitter({'type': 'status', 'data': data})
        except Exception as exc:
            log.debug('presentation status emit failed: %s', exc)

    await status('Составляю план слайдов…')
    spec = await build_outline(
        request,
        writer_model_id,
        user_prompt,
        max_slides=max_slides,
        extra_context=extra_context,
    )
    await status(f'План готов: {len(spec.slides)} слайдов')

    if enable_detail_step and len(spec.slides) > 1:

        async def detail_progress(done_n: int, total: int):
            await status(f'Дописываю слайды: {done_n}/{total}')

        spec = await detail_slides(
            request,
            writer_model_id,
            spec,
            on_progress=detail_progress,
        )

    images: dict[int, bytes] = {}
    if enable_images:
        image_targets = sum(1 for s in spec.slides if s.image_prompt)
        if image_targets:
            await status(f'Генерирую иллюстрации ({image_targets})…')

            async def img_progress(done_n: int, total: int):
                await status(f'Иллюстрации: {done_n}/{total}')

            images = await generate_slide_images(
                request,
                user,
                spec,
                concurrency=image_concurrency,
                on_progress=img_progress,
            )

    await status('Собираю .pptx…')
    # python-pptx is sync / CPU-bound → offload to a thread.
    pptx_bytes = await asyncio.to_thread(assemble_pptx, spec, images)
    await status('Презентация готова', done=True, extra={'slide_count': len(spec.slides)})

    return spec, pptx_bytes
